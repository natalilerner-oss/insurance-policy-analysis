import base64
import json
import logging
import os
import sys
import time
import uuid
from datetime import datetime, timedelta
from io import BytesIO
from typing import Any, Dict, Optional, Tuple

import azure.functions as func
from azure.storage.blob import ContentSettings, generate_blob_sas, BlobSasPermissions
from requests_toolbelt.multipart.decoder import MultipartDecoder

from src.azure_clients import (
    get_blob_service_client,
    get_blob_config_status,
    get_document_intelligence_client,
    get_openai_client,
)
from src.hebrew_utils import contains_hebrew
from src.policy_extractor import extract_policy
from src.job_store import load_job, save_job

# New imports for additional capabilities
import jwt
from jwt import InvalidTokenError
import io
import binascii
import zipfile
import psutil
import re
import threading
try:
    import resource
except ImportError:
    resource = None  # type: ignore

from pptx import Presentation
from docx import Document
import openpyxl
from azure.storage.blob import BlobClient, ContainerClient

from src.insurance_portfolio_schema import InsurancePortfolioRequest
from src.insurance_portfolio_generator import generate_insurance_portfolio

# --------------------------------------------------------------------------- 
# Configuration for new capabilities
# ---------------------------------------------------------------------------
LOG_LEVEL = os.getenv("LOG_LEVEL", "INFO").upper()
MAX_CONTENT_LENGTH = int(os.getenv("MAX_CONTENT_LENGTH", str(50 * 1024 * 1024)))
JWT_SECRET = os.environ.get("JWT_SECRET")
JWT_ALGORITHM = os.environ.get("JWT_ALGORITHM", "HS256")
VERSION = "3.1"  # Updated version
BLOB_CONNECTION_STRING = os.getenv("BLOB_CONNECTION_STRING")
BLOB_ACCOUNT_URL = os.getenv("BLOB_ACCOUNT_URL")
BLOB_CONTAINER_NAME = os.getenv("BLOB_CONTAINER_NAME")
BLOB_SAS_TOKEN = os.getenv("BLOB_SAS_TOKEN")
MAX_BLOB_FETCH_MB = int(os.getenv("MAX_BLOB_FETCH_MB", str(200)))

app = func.FunctionApp(http_auth_level=func.AuthLevel.ANONYMOUS)

# --------------------------------------------------------------------------- 
# Logging and helpers for new capabilities
# ---------------------------------------------------------------------------
_request_ctx = threading.local()

class RequestIdFilter(logging.Filter):
    def filter(self, record: logging.LogRecord) -> bool:
        record.request_id = getattr(_request_ctx, "request_id", "N/A")
        return True

logging.basicConfig(
    stream=sys.stdout,
    level=LOG_LEVEL,
    format="%(asctime)s [%(levelname)s] [insurance-func] [request_id=%(request_id)s] %(message)s",
    datefmt="%Y-%m-%d %H:%M:%S",
)

_logger = logging.getLogger("insurance-func")
_logger.addFilter(RequestIdFilter())

def get_memory_info():
    """Get current process memory usage in MB"""
    try:
        process = psutil.Process()
        mem_info = process.memory_info()
        return {
            "rss_mb": round(mem_info.rss / (1024 * 1024), 2),
            "vms_mb": round(mem_info.vms / (1024 * 1024), 2),
        }
    except (ImportError, AttributeError):
        if resource:
            try:
                usage = resource.getrusage(resource.RUSAGE_SELF)
                maxrss_mb = round(usage.ru_maxrss / 1024, 2) if sys.platform == 'darwin' else round(usage.ru_maxrss / 1024, 2)
                return {"maxrss_mb": maxrss_mb}
            except Exception:
                return {"error": "memory tracking unavailable"}
        return {"error": "psutil not found and resource module not available"}
    except Exception as e:
        return {"error": str(e)}

def _error_response(key: str, message: str, status: int = 400, *, details: Any = None) -> func.HttpResponse:
    payload = {"error": key, "message": message, "request_id": getattr(_request_ctx, "request_id", "N/A")}
    if details is not None:
        payload["details"] = details
    return func.HttpResponse(
        body=json.dumps(payload, separators=(",", ":")),
        status_code=status,
        headers={"Content-Type": "application/json", "X-Request-ID": getattr(_request_ctx, "request_id", "N/A")},
    )


def _blob_unavailable_response() -> func.HttpResponse:
    status = get_blob_config_status()
    details = {
        "configured": status.get("configured"),
        "missing": status.get("missing"),
        "required": "BLOB_CONNECTION_STRING or BLOB_ACCOUNT_URL + BLOB_SAS_TOKEN",
    }
    return _error_response("blob_unavailable", "Blob storage not configured", 503, details=details)

def _verify_jwt(req: func.HttpRequest) -> Tuple[bool, Any]:
    if not JWT_SECRET:
        return True, "jwt_not_configured"
    auth_header = req.headers.get("Authorization", "")
    if not auth_header.startswith("Bearer "):
        return False, "missing_bearer_token"
    token = auth_header.split(" ", 1)[1]
    try:
        payload = jwt.decode(token, JWT_SECRET, algorithms=[JWT_ALGORITHM])
        _logger.info("JWT validated for sub: %s", payload.get("sub"))
        return True, payload
    except InvalidTokenError as e:
        _logger.warning("Invalid JWT: %s", str(e))
        return False, f"invalid_token: {e}"

def _assign_request_id(req: func.HttpRequest) -> None:
    _request_ctx.request_id = req.headers.get("X-Request-ID") or str(uuid.uuid4())

# --------------------------------------------------------------------------- 
# Blob-backed job store for async translateFile processing.
# ---------------------------------------------------------------------------
from concurrent.futures import ThreadPoolExecutor
from threading import Lock

_job_executor = ThreadPoolExecutor(max_workers=int(os.getenv("TRANSLATEFILE_MAX_WORKERS", "4")))
_JOB_CONTAINER = os.getenv("JOB_STATE_CONTAINER", "rfp-state")
_JOB_PREFIX = os.getenv("JOB_STATE_PREFIX", "jobs")
_JOB_RETENTION_SECONDS = int(os.getenv("JOB_RETENTION_SECONDS", "3600"))
_INMEM_FALLBACK = not (os.getenv("BLOB_CONNECTION_STRING") or os.getenv("BLOB_ACCOUNT_URL"))
_fallback_jobs_lock = Lock()
_fallback_jobs: dict[str, dict[str, Any]] = {}

def _job_blob_name(job_id: str) -> str:
    return f"{_JOB_PREFIX.rstrip('/')}/{job_id}.json"

def _get_job_blob_client(job_id: str) -> BlobClient:
    if BLOB_CONNECTION_STRING:
        return BlobClient.from_connection_string(BLOB_CONNECTION_STRING, container_name=_JOB_CONTAINER, blob_name=_job_blob_name(job_id))
    if BLOB_ACCOUNT_URL:
        base = BLOB_ACCOUNT_URL.rstrip('/')
        sas = ("?" + BLOB_SAS_TOKEN.lstrip("?")) if BLOB_SAS_TOKEN else ""
        url = f"{base}/{_JOB_CONTAINER}/{_job_blob_name(job_id)}{sas}"
        return BlobClient.from_blob_url(url)
    raise RuntimeError("blob_configuration_missing: Need BLOB_CONNECTION_STRING or BLOB_ACCOUNT_URL for job state")

def _ensure_job_container() -> None:
    if _INMEM_FALLBACK:
        return
    try:
        if BLOB_CONNECTION_STRING:
            cc = ContainerClient.from_connection_string(BLOB_CONNECTION_STRING, container_name=_JOB_CONTAINER)
        else:
            base = BLOB_ACCOUNT_URL.rstrip('/')
            sas = ("?" + BLOB_SAS_TOKEN.lstrip("?")) if BLOB_SAS_TOKEN else ""
            cc = ContainerClient.from_container_url(f"{base}/{_JOB_CONTAINER}{sas}")
        try:
            cc.create_container()
        except Exception as e:
            if "ContainerAlreadyExists" not in str(e):
                _logger.debug("Container create non-fatal error: %s", e)
    except Exception as e:
        _logger.warning("Could not ensure job container exists: %s", e)

def _create_job(job_id: str, request_id: str, original_url: str | None = None) -> bool:
    meta = {"status": "pending", "createdAt": time.time(), "request_id": request_id}
    if original_url:
        meta["originalUrl"] = original_url
    if _INMEM_FALLBACK:
        with _fallback_jobs_lock:
            _fallback_jobs[job_id] = meta
        return True
    _ensure_job_container()
    try:
        bc = _get_job_blob_client(job_id)
        bc.upload_blob(json.dumps(meta).encode("utf-8"), overwrite=True, content_type="application/json")
        try:
            downloaded = bc.download_blob(max_concurrency=1).readall()
            if not downloaded:
                _logger.error("Job blob empty after upload job_id=%s", job_id)
                return False
        except Exception as vr:
            _logger.error("Verification read failed job_id=%s err=%s", job_id, vr)
            return False
        return True
    except Exception as e:
        _logger.error("Failed to create job blob job_id=%s error=%s", job_id, e)
        return False

def _set_job_result(job_id: str, result: dict[str, Any], *, failed: bool = False) -> None:
    if _INMEM_FALLBACK:
        with _fallback_jobs_lock:
            if job_id in _fallback_jobs:
                update_payload = {
                    "status": "failed" if failed else "completed",
                    "finishedAt": time.time(),
                    "result": result,
                }
                if not failed and isinstance(result, dict):
                    if "translatedBlobUrl" in result:
                        update_payload["translatedBlobUrl"] = result["translatedBlobUrl"]
                    if "translatedDownloadUrl" in result:
                        update_payload["translatedDownloadUrl"] = result["translatedDownloadUrl"]
                _fallback_jobs[job_id].update(update_payload)
        return
    try:
        bc = _get_job_blob_client(job_id)
        try:
            existing_bytes = bc.download_blob().readall()
            existing = json.loads(existing_bytes.decode("utf-8")) if existing_bytes else {}
        except Exception:
            existing = {}
        existing.update({"status": "failed" if failed else "completed", "finishedAt": time.time(), "result": result})
        if not failed and isinstance(result, dict):
            if "translatedBlobUrl" in result:
                existing["translatedBlobUrl"] = result["translatedBlobUrl"]
            if "translatedDownloadUrl" in result:
                existing["translatedDownloadUrl"] = result["translatedDownloadUrl"]
        bc.upload_blob(json.dumps(existing).encode("utf-8"), overwrite=True, content_type="application/json")
    except Exception as e:
        _logger.error("Failed to set job result job_id=%s error=%s", job_id, e)
        with _fallback_jobs_lock:
            if job_id in _fallback_jobs:
                update_payload = {
                    "status": "failed" if failed else "completed",
                    "finishedAt": time.time(),
                    "result": result,
                }
                if not failed and isinstance(result, dict):
                    if "translatedBlobUrl" in result:
                        update_payload["translatedBlobUrl"] = result["translatedBlobUrl"]
                    if "translatedDownloadUrl" in result:
                        update_payload["translatedDownloadUrl"] = result["translatedDownloadUrl"]
                _fallback_jobs[job_id].update(update_payload)

def _get_job(job_id: str) -> dict[str, Any] | None:
    if _INMEM_FALLBACK:
        with _fallback_jobs_lock:
            return _fallback_jobs.get(job_id)
    try:
        bc = _get_job_blob_client(job_id)
        data = bc.download_blob().readall()
        if not data:
            return None
        return json.loads(data.decode("utf-8"))
    except Exception as e:
        if "BlobNotFound" in str(e):
            return None
        _logger.error("Failed to get job job_id=%s error=%s", job_id, e)
        return None

def _purge_expired_jobs(max_age_seconds: int | None = None) -> None:
    retention = max_age_seconds or _JOB_RETENTION_SECONDS
    if _INMEM_FALLBACK:
        cutoff = time.time() - retention
        with _fallback_jobs_lock:
            expired = [jid for jid, meta in _fallback_jobs.items() if meta.get("finishedAt") and meta["finishedAt"] < cutoff]
            for jid in expired:
                _fallback_jobs.pop(jid, None)
        return
    try:
        if not (BLOB_CONNECTION_STRING or BLOB_ACCOUNT_URL):
            return
        from azure.storage.blob import ContainerClient
        if BLOB_CONNECTION_STRING:
            cc = ContainerClient.from_connection_string(BLOB_CONNECTION_STRING, container_name=_JOB_CONTAINER)
        else:
            base = BLOB_ACCOUNT_URL.rstrip('/')
            sas = ("?" + BLOB_SAS_TOKEN.lstrip("?")) if BLOB_SAS_TOKEN else ""
            cc = ContainerClient.from_container_url(f"{base}/{_JOB_CONTAINER}{sas}")
        cutoff = time.time() - retention
        for blob in cc.list_blobs(name_starts_with=_JOB_PREFIX.rstrip('/') + "/"):
            try:
                bc = cc.get_blob_client(blob)
                raw = bc.download_blob().readall()
                meta = json.loads(raw.decode("utf-8")) if raw else {}
                finished = meta.get("finishedAt")
                if finished and finished < cutoff:
                    bc.delete_blob()
            except Exception as ie:
                _logger.warning("Failed purging job blob=%s err=%s", blob.name, ie)
    except Exception as e:
        _logger.warning("Job purge skipped due to error: %s", e)

# --------------------------------------------------------------------------- 
# Extraction Logic (Office docs)
# ---------------------------------------------------------------------------

def extract_text_from_pptx_bytes(file_bytes: bytes) -> str:
    _logger.info("Starting PPTX text extraction, bytes=%d, mem=%s", len(file_bytes), get_memory_info())
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        _logger.info("PPTX loaded, slides=%d, mem=%s", len(prs.slides), get_memory_info())
        
        slides_text = []
        for i, slide in enumerate(prs.slides, start=1):
            slide_text = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text and shape.text.strip()]
            if slide_text:
                slides_text.append(f"--- Slide {i} ---\n" + "\n".join(slide_text))
        
        result = "\n\n".join(slides_text)
        _logger.info("PPTX extraction complete, output_length=%d, mem=%s", len(result), get_memory_info())
        return result
    except Exception as e:
        _logger.error("Error in PPTX extraction: %s, mem=%s", str(e), get_memory_info(), exc_info=True)
        raise

def extract_text_from_docx_bytes(file_bytes: bytes) -> str:
    _logger.info("Starting DOCX text extraction, bytes=%d, mem=%s", len(file_bytes), get_memory_info())
    try:
        doc = Document(io.BytesIO(file_bytes))
        _logger.info("DOCX loaded, paragraphs=%d, tables=%d, mem=%s", len(doc.paragraphs), len(doc.tables), get_memory_info())
        
        parts = []
        para_texts = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
        if para_texts:
            parts.append("--- Paragraphs ---\n" + "\n".join(para_texts))
        
        table_texts = []
        for t_index, table in enumerate(doc.tables, start=1):
            cells = ["\t".join(cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()) for row in table.rows]
            if cells:
                table_texts.append(f"--- Table {t_index} ---\n" + "\n".join(filter(None, cells)))
        if table_texts:
            parts.append("\n\n".join(table_texts))
        
        result = "\n\n".join(parts).strip()
        _logger.info("DOCX extraction complete, output_length=%d, mem=%s", len(result), get_memory_info())
        return result
    except Exception as e:
        _logger.error("Error in DOCX extraction: %s, mem=%s", str(e), get_memory_info(), exc_info=True)
        raise

def extract_text_from_xlsx_bytes(file_bytes: bytes) -> str:
    _logger.info("Starting XLSX text extraction, bytes=%d, mem=%s", len(file_bytes), get_memory_info())
    try:
        wb = openpyxl.load_workbook(filename=io.BytesIO(file_bytes), data_only=True)
        _logger.info("XLSX loaded, worksheets=%d, mem=%s", len(wb.worksheets), get_memory_info())
        
        sheets_text = []
        for sheet in wb.worksheets:
            rows_text = []
            for row in sheet.iter_rows(values_only=True):
                values = [("" if v is None else str(v)).strip() for v in row]
                if any(values):
                    rows_text.append("\t".join(values))
            if rows_text:
                sheets_text.append(f"--- Sheet: {sheet.title} ---\n" + "\n".join(rows_text))
        
        result = "\n\n".join(sheets_text).strip()
        _logger.info("XLSX extraction complete, output_length=%d, mem=%s", len(result), get_memory_info())
        return result
    except Exception as e:
        _logger.error("Error in XLSX extraction: %s, mem=%s", str(e), get_memory_info(), exc_info=True)
        raise

def extract_text_from_pdf_bytes(file_bytes: bytes) -> str:
    _logger.info("Starting PDF text extraction, bytes=%d, mem=%s", len(file_bytes), get_memory_info())
    try:
        from src.azure_clients import get_document_intelligence_client
        doc_client = get_document_intelligence_client()
        _logger.info("Document Intelligence client obtained, mem=%s", get_memory_info())
        
        # Analyze document for layout and text
        poller = doc_client.begin_analyze_document("prebuilt-read", document=file_bytes)
        result = poller.result()
        
        _logger.info("PDF analyzed, pages=%d, mem=%s", len(result.pages), get_memory_info())
        
        pages_text = []
        for i, page in enumerate(result.pages, start=1):
            page_lines = []
            for line in page.lines:
                page_lines.append(line.content)
            if page_lines:
                pages_text.append(f"--- Page {i} ---\n" + "\n".join(page_lines))
        
        result_text = "\n\n".join(pages_text)
        _logger.info("PDF extraction complete, output_length=%d, mem=%s", len(result_text), get_memory_info())
        return result_text
    except Exception as e:
        _logger.error("Error in PDF extraction: %s, mem=%s", str(e), get_memory_info(), exc_info=True)
        raise

def _validate_or_repair_zip_bytes(b: bytes) -> bytes | None:
    if zipfile.is_zipfile(io.BytesIO(b)):
        return b
    pk_header = b"PK\x03\x04"
    first_pk = b.find(pk_header)
    if first_pk > 0:
        repaired_bytes = b[first_pk:]
        if zipfile.is_zipfile(io.BytesIO(repaired_bytes)):
            _logger.info("Repaired ZIP file by trimming %d leading bytes.", first_pk)
            return repaired_bytes
    _logger.warning("Could not validate or repair ZIP file.")
    return None

def _download_blob_bytes(*, blob_url: str | None = None, container: str | None = None, blob_name: str | None = None) -> tuple[str, bytes]:
    if not blob_url and not (container and blob_name):
        raise ValueError("Must supply either blob_url or container+blob_name.")
    if blob_url:
        from urllib.parse import urlparse, unquote, parse_qs
        parsed = urlparse(blob_url)
        qs = parse_qs(parsed.query)
        has_sig = any(k.lower() == "sig" for k in qs.keys()) or "sig" in parsed.query.lower()
        filename = unquote(parsed.path.rstrip("/").split("/")[-1])
        if not has_sig and BLOB_CONNECTION_STRING and BLOB_CONTAINER_NAME:
            path_parts = [p for p in parsed.path.split("/") if p]
            if len(path_parts) >= 2 and path_parts[0] == BLOB_CONTAINER_NAME:
                blob_name_fallback = "/".join(path_parts[1:])
                try:
                    bc = BlobClient.from_connection_string(BLOB_CONNECTION_STRING, container_name=BLOB_CONTAINER_NAME, blob_name=blob_name_fallback)
                    stream = bc.download_blob()
                    data = stream.readall()
                except Exception as e:
                    _logger.error("Fallback connection string download failed: %s", e, exc_info=True)
                    raise ValueError(f"blob_download_failed: {e}")
            else:
                _logger.info("Direct URL provided without SAS; attempting anonymous download: %s", blob_url)
                try:
                    bc = BlobClient.from_blob_url(blob_url)
                    stream = bc.download_blob()
                    data = stream.readall()
                except Exception as e:
                    msg = str(e)
                    _logger.error("Failed direct blob_url download (no SAS): %s", msg, exc_info=True)
                    if "NoAuthenticationInformation" in msg or "AuthenticationFailed" in msg:
                        raise ValueError(f"blob_auth_failed: {e}")
                    raise ValueError(f"blob_download_failed: {e}")
        else:
            _logger.info("Downloading blob via direct URL: %s (has_sig=%s)", blob_url, has_sig)
            try:
                bc = BlobClient.from_blob_url(blob_url)
                stream = bc.download_blob()
                data = stream.readall()
            except Exception as e:
                msg = str(e)
                _logger.error("Failed direct blob_url download: %s", msg, exc_info=True)
                if "NoAuthenticationInformation" in msg or "AuthenticationFailed" in msg:
                    raise ValueError(f"blob_auth_failed: {e}")
                raise ValueError(f"blob_download_failed: {e}")
    else:
        if not container:
            container = BLOB_CONTAINER_NAME
        if not container or not blob_name:
            raise ValueError("container_and_blob_name_required")
        _logger.info("Downloading blob via container/blob_name: %s/%s", container, blob_name)
        try:
            if BLOB_CONNECTION_STRING:
                bc = BlobClient.from_connection_string(BLOB_CONNECTION_STRING, container_name=container, blob_name=blob_name)
            elif BLOB_ACCOUNT_URL:
                url = f"{BLOB_ACCOUNT_URL.rstrip('/')}/{container}/{blob_name}"
                if BLOB_SAS_TOKEN:
                    url = url + ("?" + BLOB_SAS_TOKEN.lstrip("?"))
                bc = BlobClient.from_blob_url(url)
            else:
                raise ValueError("blob_configuration_missing: Need BLOB_CONNECTION_STRING or BLOB_ACCOUNT_URL")
            stream = bc.download_blob()
            data = stream.readall()
            filename = blob_name
        except Exception as e:
            _logger.error("Failed container/blob download: %s", e, exc_info=True)
            raise ValueError(f"blob_download_failed: {e}")

    size = len(data)
    _logger.info("Downloaded blob '%s' size=%d bytes", filename, size)
    if size == 0:
        raise ValueError("empty_blob")
    if size > MAX_BLOB_FETCH_MB * 1024 * 1024:
        raise ValueError("blob_exceeds_max_fetch_limit")
    return filename, data


def _parse_multipart(body: bytes, content_type: str) -> Optional[Tuple[bytes, str]]:
    try:
        import cgi

        decoder = MultipartDecoder(body, content_type)
        file_parts = []
        fields: Dict[str, str] = {}

        for part in decoder.parts:
            headers = {k.decode().lower(): v.decode() for k, v in part.headers.items()}
            disposition = headers.get("content-disposition", "")
            _, params = cgi.parse_header(disposition)
            name = params.get("name")
            filename = params.get("filename")

            if filename:
                file_parts.append((name or "file", filename, part.content))
            elif name:
                fields[name] = part.text

        if "file_base64" in fields:
            filename = fields.get("filename", "document")
            return base64.b64decode(fields["file_base64"]), filename

        if file_parts:
            preferred = [p for p in file_parts if p[0] in {"file", "document", "upload"}]
            chosen = preferred[0] if preferred else file_parts[0]
            return chosen[2], chosen[1] or "upload"
        return None
    except Exception:
        return None


def _parse_request_payload(req: func.HttpRequest) -> Tuple[bytes, str]:
    content_type = req.headers.get("content-type", "").lower()
    body = req.get_body()

    if "multipart/form-data" in content_type:
        parsed = _parse_multipart(body, content_type)
        if parsed:
            return parsed

    try:
        payload = req.get_json()
    except ValueError as exc:
        raise ValueError("Invalid JSON payload") from exc

    file_base64 = payload.get("file_base64") or payload.get("data")
    if not file_base64:
        raise ValueError("file_base64 is required in JSON body")

    filename = payload.get("filename", "document")
    return base64.b64decode(file_base64), filename


def _detect_language(text: str) -> str:
    return "he" if contains_hebrew(text) else "en"


def _upload_result(job_id: str, data: Dict[str, Any]) -> Optional[str]:
    blob_service = get_blob_service_client()
    if not blob_service:
        return None

    container_name = os.environ.get("COMPLETED_JOBS_CONTAINER", "policy-extractions")
    container = blob_service.get_container_client(container_name)
    try:
        container.create_container()
    except Exception:
        pass

    blob_name = f"{job_id}.json"
    blob_client = container.get_blob_client(blob_name)
    blob_client.upload_blob(
        json.dumps(data, ensure_ascii=False),
        overwrite=True,
        content_settings=ContentSettings(content_type="application/json"),
    )

    sas = generate_blob_sas(
        account_name=blob_client.account_name,
        container_name=container_name,
        blob_name=blob_name,
        account_key=blob_service.credential.account_key,
        permission=BlobSasPermissions(read=True),
        expiry=datetime.utcnow() + timedelta(days=1),
    )
    return f"{blob_client.url}?{sas}"


def _run_extraction(job_id: str, content: bytes) -> None:
    try:
        doc_client = get_document_intelligence_client()
        openai_client = get_openai_client()
        deployment = os.environ.get("AZURE_OPENAI_DEPLOYMENT", "gpt-4.1")

        sample_text, _ = doc_client.begin_analyze_document("prebuilt-document", content).result(), None
        source_language = _detect_language(sample_text.content if sample_text else "")

        result = extract_policy(
            doc_client=doc_client,
            openai_client=openai_client,
            deployment=deployment,
            content=content,
            source_language=source_language,
        )

        download_url = _upload_result(job_id, result)
        try:
            save_job(
                job_id,
                {
                    "jobId": job_id,
                    "status": "completed",
                    "result": result,
                    "download_url": download_url,
                },
            )
        except Exception as save_exc:
            _logger.error("Failed to persist job result: %s", save_exc)
    except Exception as exc:
        try:
            save_job(
                job_id,
                {
                    "jobId": job_id,
                    "status": "failed",
                    "error": str(exc),
                },
            )
        except Exception as save_exc:
            _logger.error("Failed to persist job failure: %s", save_exc)


@app.route(route="extract_policy")
def extract_policy_sync(req: func.HttpRequest) -> func.HttpResponse:
    logging.info("Processing extract_policy request")
    try:
        content, _ = _parse_request_payload(req)
        doc_client = get_document_intelligence_client()
        openai_client = get_openai_client()
        deployment = os.environ.get("AZURE_OPENAI_DEPLOYMENT", "gpt-4.1")

        preview = doc_client.begin_analyze_document("prebuilt-document", content).result()
        source_language = _detect_language(preview.content if preview else "")

        result = extract_policy(
            doc_client=doc_client,
            openai_client=openai_client,
            deployment=deployment,
            content=content,
            source_language=source_language,
        )

        return func.HttpResponse(
            json.dumps(result, ensure_ascii=False),
            mimetype="application/json",
            status_code=200,
        )
    except Exception as exc:
        return func.HttpResponse(str(exc), status_code=400)


@app.route(route="extract_policy_async")
def extract_policy_async(req: func.HttpRequest) -> func.HttpResponse:
    logging.info("Processing extract_policy_async request")
    try:
        content, _ = _parse_request_payload(req)
    except Exception as exc:
        return func.HttpResponse(str(exc), status_code=400)

    job_id = str(uuid.uuid4())
    if not get_blob_config_status().get("configured"):
        return _blob_unavailable_response()
    try:
        save_job(
            job_id,
            {
                "jobId": job_id,
                "status": "running",
                "created_at": datetime.utcnow().isoformat(),
            },
        )
    except Exception as exc:
        return func.HttpResponse(str(exc), status_code=500)

    import threading

    thread = threading.Thread(target=_run_extraction, args=(job_id, content), daemon=True)
    thread.start()

    status_url = req.url.replace("/extract_policy_async", f"/extract_policy/status/{job_id}")
    response = {"jobId": job_id, "pollingUrl": status_url}
    return func.HttpResponse(json.dumps(response), mimetype="application/json", status_code=202)


@app.route(route="extract_policy/status/{jobId}")
def extract_policy_status(req: func.HttpRequest) -> func.HttpResponse:
    job_id = req.route_params.get("jobId")
    if not job_id:
        return func.HttpResponse("jobId is required", status_code=400)

    if not get_blob_config_status().get("configured"):
        return _blob_unavailable_response()
    try:
        job = load_job(job_id)
    except Exception as exc:
        return func.HttpResponse(str(exc), status_code=500)

    if not job:
        return func.HttpResponse("Job not found", status_code=404)

    return func.HttpResponse(json.dumps(job, ensure_ascii=False), mimetype="application/json")


# --------------------------------------------------------------------------- 
# New routes for Office document extraction and translation
# ---------------------------------------------------------------------------

@app.route(route="extract", methods=["POST"])
def extract(req: func.HttpRequest) -> func.HttpResponse:
    _assign_request_id(req)
    _logger.info("Extract route triggered")
    ok, detail = _verify_jwt(req)
    if not ok:
        return _error_response("unauthorized", "Authorization failed", 401, details=detail)

    filename: str | None = None
    file_bytes: bytes | None = None

    raw_body = req.get_body() or b""
    if raw_body:
        try:
            body_json = json.loads(raw_body.decode("utf-8")) if raw_body.strip().startswith(b"{") else None
        except Exception:
            body_json = None
    else:
        body_json = None

    if body_json:
        blob_url = body_json.get("file-url") or body_json.get("blob_url")
        container = body_json.get("container")
        blob_name = body_json.get("blob_name")
        if blob_url or blob_name:
            try:
                filename, file_bytes = _download_blob_bytes(blob_url=blob_url, container=container, blob_name=blob_name)
            except ValueError as e:
                detail_msg = str(e)
                if detail_msg.startswith("blob_auth_failed"):
                    return _error_response("unauthorized_blob", "Blob access unauthorized", 401, details=detail_msg)
                err_key = "blob_download_failed" if detail_msg.startswith("blob_download_failed") else "bad_request"
                return _error_response(err_key, "Failed to download blob", 400, details=detail_msg)

    if not file_bytes:
        # Parse multipart form using req.files if available
        if not getattr(req, "files", None) or "file" not in req.files:
            return _error_response("bad_request", "No 'file' field provided", 400, details={"available_form_fields": list(getattr(req, "form", {}).keys())})
        file = req.files["file"]
        filename = file.filename
        file_bytes = file.read()

    assert filename is not None and file_bytes is not None
    if len(file_bytes) > MAX_CONTENT_LENGTH:
        return _error_response("payload_too_large", f"File size exceeds limit of {MAX_CONTENT_LENGTH} bytes", 413)

    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""
    if ext not in {"pptx", "docx", "xlsx", "pdf"}:
        return _error_response("unsupported_media_type", "Only .pptx, .docx, .xlsx, .pdf files are supported", 415)

    if not file_bytes.startswith(b"PK\x03\x04"):
        try:
            decoded = base64.b64decode(file_bytes.strip(), validate=True)
            if decoded.startswith(b"PK\x03\x04"):
                file_bytes = decoded
        except (binascii.Error, ValueError):
            pass

    validated = _validate_or_repair_zip_bytes(file_bytes)
    if not validated and ext != "pdf":  # PDFs don't need ZIP validation
        return _error_response("bad_request", "File is not a valid Office document or is corrupted", 400)
    if validated:
        file_bytes = validated

    try:
        extraction_map = {
            "pptx": extract_text_from_pptx_bytes,
            "docx": extract_text_from_docx_bytes,
            "xlsx": extract_text_from_xlsx_bytes,
            "pdf": extract_text_from_pdf_bytes,
        }
        content = extraction_map[ext](file_bytes)
        return func.HttpResponse(content, status_code=200, headers={"Content-Type": "text/plain; charset=utf-8", "X-Request-ID": _request_ctx.request_id})
    except MemoryError:
        return _error_response("out_of_memory", "Insufficient memory to process the file", 507)
    except Exception as e:
        _logger.error("Extraction failed: %s", e, exc_info=True)
        return _error_response("failed_to_process", "Unexpected error during extraction", 500, details=str(e))


@app.route(route="translateText", methods=["POST"])
def translate_text_route(req: func.HttpRequest) -> func.HttpResponse:
    _assign_request_id(req)
    ok, detail = _verify_jwt(req)
    if not ok:
        return _error_response("unauthorized", "Authorization failed", 401, details=detail)
    try:
        body = req.get_json()
    except ValueError:
        return _error_response("bad_request", "Invalid JSON payload", 400)
    text = body.get("text")
    if not text:
        return _error_response("bad_request", "'text' is required", 400)
    to_lang = body.get("to") or os.getenv("DEFAULT_TARGET_LANG", "fr")
    from_lang = body.get("from") or os.getenv("DEFAULT_SOURCE_LANG", "en")
    try:
        translated = translate_text(text=text, to_lang=to_lang, from_lang=from_lang)
    except TranslationError as e:
        return _error_response("translation_failed", str(e), 500)
    return func.HttpResponse(
        json.dumps({"translatedText": translated, "to": to_lang, "from": from_lang, "request_id": _request_ctx.request_id}),
        status_code=200,
        mimetype="application/json",
        headers={"X-Request-ID": _request_ctx.request_id},
    )


@app.route(route="translateFile", methods=["POST"])
def translate_file_route(req: func.HttpRequest) -> func.HttpResponse:
    _assign_request_id(req)
    _logger.info("translateFile route triggered method=%s url=%s", req.method, req.url)
    
    ok, detail = _verify_jwt(req)
    if not ok:
        return _error_response("unauthorized", "Authorization failed", 401, details=detail)
    
    try:
        body = req.get_json()
        _logger.info("Request body parsed successfully")
    except ValueError as e:
        _logger.error("Invalid JSON payload: %s", e, exc_info=True)
        return _error_response("bad_request", "Invalid JSON payload", 400)
    
    file_url = body.get("fileUrl")
    if not file_url:
        _logger.error("Missing 'fileUrl' in request body")
        return _error_response("bad_request", "'fileUrl' is required", 400)
    
    to_lang = body.get("to") or os.getenv("DEFAULT_TARGET_LANG", "fr")
    from_lang = body.get("from") or os.getenv("DEFAULT_SOURCE_LANG", "en")
    correlation_id = str(uuid.uuid4())
    
    _logger.info("Starting file translation: correlation_id=%s, file_url=%s, from=%s, to=%s", 
                 correlation_id, file_url[:100], from_lang, to_lang)
    
    _purge_expired_jobs()
    job_id = correlation_id
    created_ok = _create_job(job_id, _request_ctx.request_id, original_url=file_url)
    if not created_ok:
        return _error_response("job_init_failed", "Unable to persist initial job state", 500, details={"jobId": job_id})

    def _worker():
        try:
            result = translate_remote_file(file_url=file_url, to_lang=to_lang, from_lang=from_lang, correlation_id=correlation_id)
            result["request_id"] = _request_ctx.request_id
            if "error" in result:
                _logger.error("Async translation failed job_id=%s error=%s", job_id, result.get("error"))
                _set_job_result(job_id, result, failed=True)
            else:
                _logger.info("Async translation completed job_id=%s", job_id)
                _set_job_result(job_id, result, failed=False)
        except MemoryError:
            err = {"error": "out_of_memory", "message": "Insufficient memory", "correlationId": correlation_id, "memory": get_memory_info()}
            _set_job_result(job_id, err, failed=True)
        except TimeoutError:
            err = {"error": "timeout", "message": "Processing exceeded time limit", "correlationId": correlation_id}
            _set_job_result(job_id, err, failed=True)
        except Exception as e:
            err = {"error": "remote_translation_failed", "message": str(e), "correlationId": correlation_id, "error_type": type(e).__name__, "memory": get_memory_info()}
            _set_job_result(job_id, err, failed=True)

    parent_request_id = getattr(_request_ctx, "request_id", None)

    def _worker_wrapper():
        if parent_request_id:
            _request_ctx.request_id = parent_request_id
        try:
            _worker()
        finally:
            if hasattr(_request_ctx, "request_id"):
                delattr(_request_ctx, "request_id")

    _job_executor.submit(_worker_wrapper)

    polling_url = f"{req.url.rstrip('/')}/status/{job_id}"
    accepted_body = {
        "status": "accepted",
        "jobId": job_id,
        "pollingUrl": polling_url,
        "request_id": _request_ctx.request_id,
    }
    return func.HttpResponse(
        json.dumps(accepted_body),
        status_code=202,
        mimetype="application/json",
        headers={
            "X-Request-ID": _request_ctx.request_id,
            "Location": polling_url,
            "Retry-After": os.getenv("TRANSLATEFILE_RETRY_AFTER", "3"),
        },
    )


@app.route(route="translateFile/status/{job_id}", methods=["GET"])
def translate_file_status_route(req: func.HttpRequest) -> func.HttpResponse:
    _assign_request_id(req)
    job_id = getattr(req, "route_params", {}).get("job_id")
    if not job_id:
        return _error_response("bad_request", "Missing job_id route parameter", 400)
    job = _get_job(job_id)
    if not job:
        return _error_response("not_found", f"Job '{job_id}' not found", 404)
    status = job.get("status")
    if status == "pending":
        body = {
            "status": "pending",
            "jobId": job_id,
            "request_id": _request_ctx.request_id,
        }
        return func.HttpResponse(
            json.dumps(body),
            status_code=202,
            mimetype="application/json",
            headers={"X-Request-ID": _request_ctx.request_id, "Retry-After": os.getenv("TRANSLATEFILE_RETRY_AFTER", "3")},
        )
    result = job.get("result", {})
    body = {
        "status": status,
        "jobId": job_id,
        "result": result,
        "request_id": _request_ctx.request_id,
    }
    return func.HttpResponse(
        json.dumps(body),
        status_code=200,
        mimetype="application/json",
        headers={"X-Request-ID": _request_ctx.request_id},
    )


@app.route(route="health", methods=["GET"])
def health(req: func.HttpRequest) -> func.HttpResponse:
    _assign_request_id(req)
    body = {
        "status": "ok",
        "version": VERSION,
        "jwtEnabled": bool(JWT_SECRET),
        "blobMode": bool(BLOB_CONNECTION_STRING or BLOB_ACCOUNT_URL),
        "maxUploadMB": round(MAX_CONTENT_LENGTH / (1024 * 1024), 2),
        "maxBlobFetchMB": MAX_BLOB_FETCH_MB,
        "supportedFormats": ["pptx", "docx", "xlsx", "pdf"],
        "capabilities": [
            "insurance_policy_extraction",
            "office_document_text_extraction",
            "text_translation",
            "file_translation",
            "insurance_portfolio_excel_generation",
            "policy_comparison"
        ],
        "memory": get_memory_info(),
        "request_id": _request_ctx.request_id,
    }
    return func.HttpResponse(json.dumps(body), status_code=200, mimetype="application/json", headers={"X-Request-ID": _request_ctx.request_id})


@app.route(route="compare_policies", methods=["POST"])
def compare_policies(req: func.HttpRequest) -> func.HttpResponse:
    _assign_request_id(req)
    _logger.info("compare_policies route triggered")
    ok, detail = _verify_jwt(req)
    if not ok:
        return _error_response("unauthorized", "Authorization failed", 401, details=detail)

    try:
        body = req.get_json()
    except ValueError as exc:
        return _error_response("bad_request", f"Invalid JSON payload: {str(exc)}", 400)

    policies = body.get("policies") if isinstance(body, dict) else body
    if not isinstance(policies, list) or not policies:
        return _error_response("bad_request", "Request must include a non-empty policies array", 400)

    def policy_label(policy: Dict[str, Any], index: int) -> str:
        policy_number = policy.get("policy_number") or f"policy_{index + 1}"
        carrier = policy.get("carrier", {}).get("name")
        return f"{policy_number} ({carrier})" if carrier else policy_number

    def coverage_premium(coverage: Dict[str, Any]) -> float:
        premium = coverage.get("premium")
        if isinstance(premium, dict):
            return float(premium.get("final_monthly") or premium.get("base_monthly") or 0)
        if premium is None:
            return 0.0
        return float(premium)

    policy_labels = []
    summary = []
    coverage_index: Dict[str, Dict[str, float]] = {}

    for idx, policy in enumerate(policies):
        if not isinstance(policy, dict):
            continue
        label = policy_label(policy, idx)
        policy_labels.append(label)
        total_premium = policy.get("total_monthly_premium") or 0
        summary.append(
            {
                "policy": label,
                "policy_number": policy.get("policy_number"),
                "carrier": policy.get("carrier", {}).get("name"),
                "total_monthly_premium": float(total_premium or 0),
            }
        )

        for coverage in policy.get("coverages", []) or []:
            if not isinstance(coverage, dict):
                continue
            coverage_name = coverage.get("type") or coverage.get("product_name") or "Coverage"
            coverage_index.setdefault(coverage_name, {})[label] = coverage_premium(coverage)

    comparison_rows = []
    for coverage_name, by_policy in coverage_index.items():
        row = {"coverage": coverage_name}
        for label in policy_labels:
            row[label] = by_policy.get(label, 0)
        comparison_rows.append(row)

    response = {
        "policies": summary,
        "columns": ["coverage"] + policy_labels,
        "rows": comparison_rows,
        "request_id": _request_ctx.request_id,
    }
    return func.HttpResponse(
        json.dumps(response, ensure_ascii=False),
        status_code=200,
        mimetype="application/json",
        headers={"X-Request-ID": _request_ctx.request_id},
    )


# --------------------------------------------------------------------------- 
# Insurance Portfolio Excel Generation
# ---------------------------------------------------------------------------

@app.route(route="generate_insurance_portfolio", methods=["POST"])
def generate_insurance_portfolio_endpoint(req: func.HttpRequest) -> func.HttpResponse:
    _assign_request_id(req)
    _logger.info("generate_insurance_portfolio route triggered")
    ok, detail = _verify_jwt(req)
    if not ok:
        return _error_response("unauthorized", "Authorization failed", 401, details=detail)

    try:
        body = req.get_json()
        data = InsurancePortfolioRequest(**body)
    except ValueError as e:
        return _error_response("bad_request", f"Invalid request data: {str(e)}", 400)

    try:
        excel_bytes = generate_insurance_portfolio(data)

        # Calculate totals for summary (needed for both paths)
        total_premium = sum(
            sum(c.premium for c in p.coverages) if p.coverages else (p.premium or 0)
            for p in data.insurance_products
        )
        summary = {
            "family_name": data.family_name,
            "report_date": str(data.report_date),
            "products_count": len(data.insurance_products),
            "total_monthly_premium": float(total_premium),
            "family_members_count": len(data.family_members)
        }

        # Upload to blob storage
        filename = f"תיק_ביטוח_משפחת_{data.family_name}_{data.report_date}.xlsx"
        blob_service = get_blob_service_client()
        download_url = None

        if blob_service:
            try:
                container_name = os.environ.get("EXCEL_REPORTS_CONTAINER", "excel-reports")
                container = blob_service.get_container_client(container_name)
                try:
                    container.create_container()
                except Exception:
                    pass

                blob_name = f"{data.family_name}_{data.report_date}.xlsx"
                blob_client = container.get_blob_client(blob_name)
                blob_client.upload_blob(
                    excel_bytes.getvalue(),
                    overwrite=True,
                    content_settings=ContentSettings(content_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"),
                )

                # Generate SAS URL
                sas = generate_blob_sas(
                    account_name=blob_client.account_name,
                    container_name=container_name,
                    blob_name=blob_name,
                    account_key=blob_service.credential.account_key,
                    permission=BlobSasPermissions(read=True),
                    expiry=datetime.utcnow() + timedelta(days=7),  # Longer expiry for reports
                )
                download_url = f"{blob_client.url}?{sas}"
            except Exception as e:
                _logger.error("Failed to upload report to blob storage: %s", e)
                # Continue to direct download fallback

        if download_url:
            response_data = {
                "success": True,
                "filename": filename,
                "downloadUrl": download_url,
                "summary": summary,
                "request_id": _request_ctx.request_id
            }

            return func.HttpResponse(
                json.dumps(response_data, ensure_ascii=False),
                status_code=200,
                mimetype="application/json",
                headers={"X-Request-ID": _request_ctx.request_id}
            )
        else:
            # Fallback: Return file directly
            return func.HttpResponse(
                excel_bytes.getvalue(),
                mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                status_code=200,
                headers={
                    "Content-Disposition": f"attachment; filename={filename}",
                    "X-Request-ID": _request_ctx.request_id
                }
            )

    except Exception as e:
        _logger.error("Failed to generate insurance portfolio: %s", e, exc_info=True)
        return _error_response("failed_to_generate", f"Excel generation failed: {str(e)}", 500)
