import base64
import binascii
import cgi
import io
import json
import logging
import os
import threading
import uuid
import zipfile
from datetime import datetime, timedelta
from typing import Any, Dict, Optional, Tuple

import azure.functions as func
from azure.storage.blob import ContentSettings, generate_blob_sas, BlobSasPermissions, BlobClient
from requests_toolbelt.multipart.decoder import MultipartDecoder

# Office extraction imports
from pptx import Presentation
from docx import Document
import openpyxl

from src.azure_clients import (
    get_blob_service_client,
    get_blob_config_status,
    get_document_intelligence_client,
    get_openai_client,
)
from src.hebrew_utils import contains_hebrew
from src.policy_extractor import extract_policy
from src.job_store import save_job, load_job
from src.blueprints.utils import assign_request_id, error_response, verify_jwt, logger, get_memory_info, get_request_id

# Config
content_type_map = {
    "pdf": "application/pdf",
    "png": "image/png",
    "jpeg": "image/jpeg",
    "tiff": "image/tiff"
}
MAX_CONTENT_LENGTH = int(os.getenv("MAX_CONTENT_LENGTH", str(50 * 1024 * 1024)))
BLOB_CONNECTION_STRING = os.getenv("BLOB_CONNECTION_STRING")
BLOB_ACCOUNT_URL = os.getenv("BLOB_ACCOUNT_URL")
BLOB_CONTAINER_NAME = os.getenv("BLOB_CONTAINER_NAME")
BLOB_SAS_TOKEN = os.getenv("BLOB_SAS_TOKEN")

bp = func.Blueprint()

def _blob_unavailable_response() -> func.HttpResponse:
    status = get_blob_config_status()
    details = {
        "configured": status.get("configured"),
        "missing": status.get("missing"),
        "required": "BLOB_CONNECTION_STRING or BLOB_ACCOUNT_URL + BLOB_SAS_TOKEN",
    }
    return error_response("blob_unavailable", "Blob storage not configured", 503, details=details)

def _parse_multipart(body: bytes, content_type: str) -> Optional[Tuple[bytes, str]]:
    try:
        decoder = MultipartDecoder(body, content_type)
        file_parts = []
        fields: Dict[str, str] = {}

        for part in decoder.parts:
            headers = {k.decode().lower(): v.decode() for k, v in part.headers.items()}
            disposition = headers.get("content-disposition", "")
            _, params = cgi.parse_header(disposition)
            name = params.get("name")
            filename = params.get("filename")

            if filename:
                file_parts.append((name or "file", filename, part.content))
            elif name:
                fields[name] = part.text

        if "file_base64" in fields:
            filename = fields.get("filename", "document")
            return base64.b64decode(fields["file_base64"]), filename

        if file_parts:
            preferred = [p for p in file_parts if p[0] in {"file", "document", "upload"}]
            chosen = preferred[0] if preferred else file_parts[0]
            return chosen[2], chosen[1] or "upload"
        return None
    except Exception:
        return None

def _parse_request_payload(req: func.HttpRequest) -> Tuple[bytes, str]:
    content_type = req.headers.get("content-type", "").lower()
    body = req.get_body()

    if "multipart/form-data" in content_type:
        parsed = _parse_multipart(body, content_type)
        if parsed:
            return parsed

    try:
        payload = req.get_json()
    except ValueError as exc:
        raise ValueError("Invalid JSON payload") from exc

    file_base64 = payload.get("file_base64") or payload.get("data")
    if not file_base64:
        raise ValueError("file_base64 is required in JSON body")

    filename = payload.get("filename", "document")
    return base64.b64decode(file_base64), filename

def _detect_language(text: str) -> str:
    return "he" if contains_hebrew(text) else "en"

def _upload_result(job_id: str, data: Dict[str, Any]) -> Optional[str]:
    blob_service = get_blob_service_client()
    if not blob_service:
        return None

    container_name = os.environ.get("COMPLETED_JOBS_CONTAINER", "policy-extractions")
    container = blob_service.get_container_client(container_name)
    try:
        container.create_container()
    except Exception:
        pass

    blob_name = f"{job_id}.json"
    blob_client = container.get_blob_client(blob_name)
    blob_client.upload_blob(
        json.dumps(data, ensure_ascii=False),
        overwrite=True,
        content_settings=ContentSettings(content_type="application/json"),
    )

    sas = generate_blob_sas(
        account_name=blob_client.account_name,
        container_name=container_name,
        blob_name=blob_name,
        account_key=blob_service.credential.account_key,
        permission=BlobSasPermissions(read=True),
        expiry=datetime.utcnow() + timedelta(days=1),
    )
    return f"{blob_client.url}?{sas}"

def _run_extraction(job_id: str, content: bytes) -> None:
    try:
        doc_client = get_document_intelligence_client()
        openai_client = get_openai_client()
        deployment = os.environ.get("AZURE_OPENAI_DEPLOYMENT", "gpt-4.1")

        sample_text, _ = doc_client.begin_analyze_document("prebuilt-document", content).result(), None
        source_language = _detect_language(sample_text.content if sample_text else "")

        result = extract_policy(
            doc_client=doc_client,
            openai_client=openai_client,
            deployment=deployment,
            content=content,
            source_language=source_language,
        )

        download_url = _upload_result(job_id, result)
        try:
            save_job(
                job_id,
                {
                    "jobId": job_id,
                    "status": "completed",
                    "result": result,
                    "download_url": download_url,
                },
            )
        except Exception as save_exc:
            logger.error("Failed to persist job result: %s", save_exc)
    except Exception as exc:
        try:
            save_job(
                job_id,
                {
                    "jobId": job_id,
                    "status": "failed",
                    "error": str(exc),
                },
            )
        except Exception as save_exc:
             logger.error("Failed to persist job failure: %s", save_exc)

# Office Extraction functions
def extract_text_from_pptx_bytes(file_bytes: bytes) -> str:
    logger.info("Starting PPTX text extraction, bytes=%d, mem=%s", len(file_bytes), get_memory_info())
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        slides_text = []
        for i, slide in enumerate(prs.slides, start=1):
            slide_text = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text and shape.text.strip()]
            if slide_text:
                slides_text.append(f"--- Slide {i} ---\n" + "\n".join(slide_text))
        return "\n\n".join(slides_text)
    except Exception as e:
        logger.error("Error in PPTX extraction: %s", str(e), exc_info=True)
        raise

def extract_text_from_docx_bytes(file_bytes: bytes) -> str:
    logger.info("Starting DOCX text extraction, bytes=%d, mem=%s", len(file_bytes), get_memory_info())
    try:
        doc = Document(io.BytesIO(file_bytes))
        parts = []
        para_texts = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
        if para_texts:
            parts.append("--- Paragraphs ---\n" + "\n".join(para_texts))
        
        table_texts = []
        for t_index, table in enumerate(doc.tables, start=1):
            cells = ["\t".join(cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()) for row in table.rows]
            if cells:
                table_texts.append(f"--- Table {t_index} ---\n" + "\n".join(filter(None, cells)))
        if table_texts:
            parts.append("\n\n".join(table_texts))
        
        return "\n\n".join(parts).strip()
    except Exception as e:
        logger.error("Error in DOCX extraction: %s", str(e), exc_info=True)
        raise

def extract_text_from_xlsx_bytes(file_bytes: bytes) -> str:
    logger.info("Starting XLSX text extraction, bytes=%d, mem=%s", len(file_bytes), get_memory_info())
    try:
        wb = openpyxl.load_workbook(filename=io.BytesIO(file_bytes), data_only=True)
        sheets_text = []
        for sheet in wb.worksheets:
            rows_text = []
            for row in sheet.iter_rows(values_only=True):
                values = [("" if v is None else str(v)).strip() for v in row]
                if any(values):
                    rows_text.append("\t".join(values))
            if rows_text:
                sheets_text.append(f"--- Sheet: {sheet.title} ---\n" + "\n".join(rows_text))
        return "\n\n".join(sheets_text).strip()
    except Exception as e:
        logger.error("Error in XLSX extraction: %s", str(e), exc_info=True)
        raise

def extract_text_from_pdf_bytes(file_bytes: bytes) -> str:
    logger.info("Starting PDF text extraction, bytes=%d, mem=%s", len(file_bytes), get_memory_info())
    try:
        doc_client = get_document_intelligence_client()
        poller = doc_client.begin_analyze_document("prebuilt-read", document=file_bytes)
        result = poller.result()
        pages_text = []
        for i, page in enumerate(result.pages, start=1):
            page_lines = [line.content for line in page.lines]
            if page_lines:
                pages_text.append(f"--- Page {i} ---\n" + "\n".join(page_lines))
        return "\n\n".join(pages_text)
    except Exception as e:
        logger.error("Error in PDF extraction: %s", str(e), exc_info=True)
        raise

def _validate_or_repair_zip_bytes(b: bytes) -> bytes | None:
    if zipfile.is_zipfile(io.BytesIO(b)):
        return b
    pk_header = b"PK\x03\x04"
    first_pk = b.find(pk_header)
    if first_pk > 0:
        repaired_bytes = b[first_pk:]
        if zipfile.is_zipfile(io.BytesIO(repaired_bytes)):
            logger.info("Repaired ZIP file by trimming %d leading bytes.", first_pk)
            return repaired_bytes
    logger.warning("Could not validate or repair ZIP file.")
    return None

def _download_blob_bytes(*, blob_url: str | None = None, container: str | None = None, blob_name: str | None = None) -> tuple[str, bytes]:
    if not blob_url and not (container and blob_name):
        raise ValueError("Must supply either blob_url or container+blob_name.")
    try:
        if blob_url:
             bc = BlobClient.from_blob_url(blob_url)
             stream = bc.download_blob()
             return str(blob_url), stream.readall()
        else:
             if not container: container = BLOB_CONTAINER_NAME
             if BLOB_CONNECTION_STRING:
                  bc = BlobClient.from_connection_string(BLOB_CONNECTION_STRING, container_name=container, blob_name=blob_name)
             elif BLOB_ACCOUNT_URL:
                  url = f"{BLOB_ACCOUNT_URL.rstrip('/')}/{container}/{blob_name}"
                  bc = BlobClient.from_blob_url(url)
             else:
                  raise ValueError("No blob config found")
             
             stream = bc.download_blob()
             return f"{container}/{blob_name}", stream.readall()

    except Exception as e:
         raise ValueError(f"blob_download_failed: {e}")

# ROUTES

@bp.route(route="extract_policy")
def extract_policy_sync(req: func.HttpRequest) -> func.HttpResponse:
    assign_request_id(req)
    logger.info("Processing extract_policy request")
    try:
        content, _ = _parse_request_payload(req)
        doc_client = get_document_intelligence_client()
        openai_client = get_openai_client()
        deployment = os.environ.get("AZURE_OPENAI_DEPLOYMENT", "gpt-4.1")

        preview = doc_client.begin_analyze_document("prebuilt-document", content).result()
        source_language = _detect_language(preview.content if preview else "")

        result = extract_policy(
            doc_client=doc_client,
            openai_client=openai_client,
            deployment=deployment,
            content=content,
            source_language=source_language,
        )

        return func.HttpResponse(
            json.dumps(result, ensure_ascii=False),
            mimetype="application/json",
            status_code=200,
        )
    except Exception as exc:
        return func.HttpResponse(str(exc), status_code=400)

@bp.route(route="extract_policy_async")
def extract_policy_async(req: func.HttpRequest) -> func.HttpResponse:
    assign_request_id(req)
    logger.info("Processing extract_policy_async request")
    try:
        content, _ = _parse_request_payload(req)
    except Exception as exc:
        return func.HttpResponse(str(exc), status_code=400)

    job_id = str(uuid.uuid4())
    if not get_blob_config_status().get("configured"):
        return _blob_unavailable_response()
    try:
        save_job(
            job_id,
            {
                "jobId": job_id,
                "status": "running",
                "created_at": datetime.utcnow().isoformat(),
            },
        )
    except Exception as exc:
        return func.HttpResponse(str(exc), status_code=500)

    thread = threading.Thread(target=_run_extraction, args=(job_id, content), daemon=True)
    thread.start()

    status_url = req.url.replace("/extract_policy_async", f"/extract_policy/status/{job_id}")
    response = {"jobId": job_id, "pollingUrl": status_url}
    return func.HttpResponse(json.dumps(response), mimetype="application/json", status_code=202)

@bp.route(route="extract_policy/status/{jobId}")
def extract_policy_status(req: func.HttpRequest) -> func.HttpResponse:
    assign_request_id(req)
    job_id = req.route_params.get("jobId")
    if not job_id:
        return func.HttpResponse("jobId is required", status_code=400)

    if not get_blob_config_status().get("configured"):
        return _blob_unavailable_response()
    try:
        job = load_job(job_id)
    except Exception as exc:
        return func.HttpResponse(str(exc), status_code=500)

    if not job:
        return func.HttpResponse("Job not found", status_code=404)

    return func.HttpResponse(json.dumps(job, ensure_ascii=False), mimetype="application/json")


@bp.route(route="extract", methods=["POST"])
def extract(req: func.HttpRequest) -> func.HttpResponse:
    assign_request_id(req)
    logger.info("Extract route triggered")
    ok, detail = verify_jwt(req)
    if not ok:
        return error_response("unauthorized", "Authorization failed", 401, details=detail)

    filename: str | None = None
    file_bytes: bytes | None = None

    raw_body = req.get_body() or b""
    if raw_body:
        try:
            body_json = json.loads(raw_body.decode("utf-8")) if raw_body.strip().startswith(b"{") else None
        except Exception:
            body_json = None
    else:
        body_json = None

    if body_json:
        blob_url = body_json.get("file-url") or body_json.get("blob_url")
        container = body_json.get("container")
        blob_name = body_json.get("blob_name")
        if blob_url or blob_name:
            try:
                filename, file_bytes = _download_blob_bytes(blob_url=blob_url, container=container, blob_name=blob_name)
            except ValueError as e:
                detail_msg = str(e)
                if detail_msg.startswith("blob_auth_failed"):
                    return error_response("unauthorized_blob", "Blob access unauthorized", 401, details=detail_msg)
                err_key = "blob_download_failed" if detail_msg.startswith("blob_download_failed") else "bad_request"
                return error_response(err_key, "Failed to download blob", 400, details=detail_msg)

    if not file_bytes:
        # Check req.files - but req.files from function_app.py relies on standard Azure Functions handling
        if not hasattr(req, "files") or "file" not in req.files:
             # Try custom parse if needed, but assuming Azure Functions V2 supports it.
             # If req.files is missing, try our `_parse_multipart` on the body
             content_type = req.headers.get("content-type", "").lower()
             if "multipart/form-data" in content_type:
                 parsed = _parse_multipart(raw_body, content_type)
                 if parsed:
                     file_bytes, filename = parsed
        else:
             file = req.files["file"]
             filename = file.filename
             file_bytes = file.read()
    
    if not file_bytes or not filename:
         return error_response("bad_request", "No file provided", 400)

    if len(file_bytes) > MAX_CONTENT_LENGTH:
        return error_response("payload_too_large", f"File size exceeds limit of {MAX_CONTENT_LENGTH} bytes", 413)

    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""
    if ext not in {"pptx", "docx", "xlsx", "pdf"}:
        return error_response("unsupported_media_type", "Only .pptx, .docx, .xlsx, .pdf files are supported", 415)

    if not file_bytes.startswith(b"PK\x03\x04"):
        try:
            decoded = base64.b64decode(file_bytes.strip(), validate=True)
            if decoded.startswith(b"PK\x03\x04"):
                file_bytes = decoded
        except (binascii.Error, ValueError):
            pass

    validated = _validate_or_repair_zip_bytes(file_bytes)
    if not validated and ext != "pdf":  # PDFs don't need ZIP validation
        return error_response("bad_request", "File is not a valid Office document or is corrupted", 400)
    if validated:
        file_bytes = validated

    try:
        extraction_map = {
            "pptx": extract_text_from_pptx_bytes,
            "docx": extract_text_from_docx_bytes,
            "xlsx": extract_text_from_xlsx_bytes,
            "pdf": extract_text_from_pdf_bytes,
        }
        content = extraction_map[ext](file_bytes)
        return func.HttpResponse(content, status_code=200, headers={"Content-Type": "text/plain; charset=utf-8", "X-Request-ID": get_request_id()})
    except MemoryError:
        return error_response("out_of_memory", "Insufficient memory to process the file", 507)
    except Exception as e:
        logger.error("Extraction failed: %s", e, exc_info=True)
        return error_response("failed_to_process", "Unexpected error during extraction", 500, details=str(e))
